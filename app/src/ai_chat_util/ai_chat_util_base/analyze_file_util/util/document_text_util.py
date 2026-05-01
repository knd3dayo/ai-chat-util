from __future__ import annotations

import base64
import tempfile
from pathlib import Path

from docx import Document as DocxDocument
from openpyxl import load_workbook
from pdfminer.high_level import extract_text as extract_pdf_text
from pptx import Presentation


class DocumentTextUtil:
    TEXT_FILE_SUFFIXES = {
        ".txt",
        ".log",
        ".md",
        ".csv",
        ".json",
        ".yaml",
        ".yml",
        ".xml",
        ".html",
        ".htm",
        ".py",
        ".js",
        ".ts",
        ".sql",
    }

    @classmethod
    def extract_docx_text(cls, path: str | Path) -> str:
        document = DocxDocument(str(path))
        lines = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
        return "\n".join(lines)

    @classmethod
    def extract_pptx_text(cls, path: str | Path) -> str:
        presentation = Presentation(str(path))
        lines: list[str] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text and text.strip():
                    lines.append(f"[slide {slide_index}] {text.strip()}")
        return "\n".join(lines)

    @classmethod
    def extract_xlsx_text(cls, path: str | Path) -> str:
        workbook = load_workbook(path, read_only=True, data_only=True)
        try:
            lines: list[str] = []
            for sheet in workbook.worksheets:
                lines.append(f"[sheet] {sheet.title}")
                for row in sheet.iter_rows(values_only=True):
                    values = [str(value).strip() for value in row if value is not None and str(value).strip()]
                    if values:
                        lines.append("\t".join(values))
            return "\n".join(lines)
        finally:
            workbook.close()

    @classmethod
    def extract_text_from_path(cls, path: str | Path) -> str:
        resolved_path = Path(path)
        suffix = resolved_path.suffix.lower()
        if suffix in cls.TEXT_FILE_SUFFIXES:
            return resolved_path.read_text(encoding="utf-8", errors="ignore")
        if suffix == ".pdf":
            return extract_pdf_text(str(resolved_path))
        if suffix == ".docx":
            return cls.extract_docx_text(resolved_path)
        if suffix == ".pptx":
            return cls.extract_pptx_text(resolved_path)
        if suffix in {".xlsx", ".xlsm"}:
            return cls.extract_xlsx_text(resolved_path)
        raise ValueError(f"Text extraction is not supported for file type: {resolved_path.suffix or '<none>'}")

    @classmethod
    def extract_base64_to_text(cls, extension: str, base64_data: str) -> str:
        suffix = extension if extension.startswith(".") else f".{extension}"
        raw = base64.b64decode(base64_data)
        with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as handle:
            temp_path = Path(handle.name)
            handle.write(raw)
        try:
            return cls.extract_text_from_path(temp_path)
        finally:
            temp_path.unlink(missing_ok=True)