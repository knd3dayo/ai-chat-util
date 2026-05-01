from __future__ import annotations

from typing import Optional
from io import BytesIO
import os
import uuid
import tempfile
import atexit
import base64
from abc import ABC, abstractmethod

from docx import Document as WordDocument
from openpyxl import load_workbook
from pptx import Presentation

from ai_chat_util.core.common.config.runtime import AiChatUtilConfig, get_runtime_config
from ai_chat_util.core.chat.model import (
    ChatHistory, ChatRequestContext, ChatMessage,
    ChatContent, WebRequestModel
)
from ai_chat_util.core.analysis.model import FileUtilDocument
from ai_chat_util.util.analyze_file_util.office2pdf import Office2PDFUtil
from ai_chat_util.util.analyze_file_util.downloader import DownLoader
from ai_chat_util.util.analyze_file_util import pdf_util
from ai_chat_util.core.chat.abstract_chat_client import AbstractChatClient

import ai_chat_util.core.log.log_settings as log_settings
logger = log_settings.getLogger(__name__)


class FileUtilLLMMessages:
    _OFFICE_EXTRACT_MAX_SHEETS = 5
    _OFFICE_EXTRACT_MAX_ROWS_PER_SHEET = 200
    _OFFICE_EXTRACT_MAX_COLS_PER_ROW = 20

    def __init__(self, llm_client: AbstractChatClient):
        self.llm_client = llm_client

    def _get_effective_config(self) -> AiChatUtilConfig:
        config = self.llm_client.get_config()
        if config is not None:
            return config
        return get_runtime_config()

    def _get_network_download_options(self) -> tuple[bool, str | None]:
        config = self._get_effective_config()
        return config.network.requests_verify, config.network.ca_bundle

    def _extract_word_text(self, document_type: FileUtilDocument) -> str:
        document = WordDocument(BytesIO(document_type.data))
        parts: list[str] = []

        for paragraph in document.paragraphs:
            text = paragraph.text.strip()
            if text:
                parts.append(text)

        for table in document.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                if any(cells):
                    parts.append(" | ".join(cells))

        return "\n".join(parts).strip()

    def _extract_excel_text(self, document_type: FileUtilDocument) -> str:
        workbook = load_workbook(filename=BytesIO(document_type.data), read_only=True, data_only=True)
        parts: list[str] = []

        for sheet_index, sheet_name in enumerate(workbook.sheetnames, start=1):
            if sheet_index > self._OFFICE_EXTRACT_MAX_SHEETS:
                parts.append("[Truncated additional sheets]")
                break

            sheet = workbook[sheet_name]
            parts.append(f"[Sheet] {sheet_name}")
            row_count = 0
            for row in sheet.iter_rows(values_only=True):
                values = ["" if value is None else str(value).strip() for value in row[: self._OFFICE_EXTRACT_MAX_COLS_PER_ROW]]
                if not any(values):
                    continue
                parts.append(" | ".join(values))
                row_count += 1
                if row_count >= self._OFFICE_EXTRACT_MAX_ROWS_PER_SHEET:
                    parts.append("[Truncated additional rows]")
                    break

        workbook.close()
        return "\n".join(parts).strip()

    def _extract_ppt_text(self, document_type: FileUtilDocument) -> str:
        presentation = Presentation(BytesIO(document_type.data))
        parts: list[str] = []

        for slide_index, slide in enumerate(presentation.slides, start=1):
            parts.append(f"[Slide {slide_index}]")
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                normalized = text.strip() if isinstance(text, str) else ""
                if normalized:
                    parts.append(normalized)

        return "\n".join(parts).strip()

    def _extract_office_text(self, document_type: FileUtilDocument) -> str:
        if document_type.is_word():
            return self._extract_word_text(document_type)
        if document_type.is_excel():
            return self._extract_excel_text(document_type)
        if document_type.is_ppt():
            return self._extract_ppt_text(document_type)
        return ""

    def _create_office_text_fallback_content(self, document_type: FileUtilDocument) -> list["ChatContent"]:
        extracted_text = self._extract_office_text(document_type)
        if not extracted_text:
            raise RuntimeError(f"Failed to extract text from office document: {document_type.identifier}")

        explanation_content = self.create_text_content(
            text=(
                f"LibreOffice が利用できないため、Office ドキュメント {document_type.identifier} から "
                "直接テキストを抽出した内容を以下に示します。"
            )
        )
        body_content = self.create_text_content(text=extracted_text)
        return [explanation_content, body_content]

    def is_text_content(self, content: ChatContent) -> bool:
        return content.params.get("type") == "text"

    def is_image_content(self, content: ChatContent) -> bool:
        return content.params.get("type") == "image_url"

    def is_file_content(self, content: ChatContent) -> bool:
        return content.params.get("type") == "file"

    def create_text_content(self, text: str) -> "ChatContent":
        params = {"type": "text", "text": text}
        return ChatContent(params=params)

    def create_image_content(self, file_data: FileUtilDocument, detail: str) -> list["ChatContent"]:
        return self.llm_client.get_message_factory()._create_image_content_(
            file_data.identifier, file_data.data, detail
        )

    def create_image_content_from_file(self, file_path: str, detail: str) -> list["ChatContent"]:
        return self.create_image_content(FileUtilDocument.from_file(file_path), detail)

    def create_image_content_from_url(self, file_url: WebRequestModel, detail: str) -> list["ChatContent"]:
        tmpdir = tempfile.TemporaryDirectory()
        atexit.register(tmpdir.cleanup)

        requests_verify, ca_bundle = self._get_network_download_options()
        file_paths = DownLoader.download_files(
            [file_url],
            tmpdir.name,
            requests_verify=requests_verify,
            ca_bundle=ca_bundle,
        )
        document = FileUtilDocument.from_file(file_paths[0])
        return self.llm_client.get_message_factory()._create_image_content_(
            document.identifier, document.data, detail
        )

    async def create_image_content_from_url_async(self, file_url: WebRequestModel, detail: str) -> list["ChatContent"]:
        with tempfile.TemporaryDirectory() as tmp_dir:
            requests_verify, ca_bundle = self._get_network_download_options()
            file_paths = await DownLoader.download_files_async(
                [file_url],
                tmp_dir,
                requests_verify=requests_verify,
                ca_bundle=ca_bundle,
            )
            document = FileUtilDocument.from_file(file_paths[0])
            return self.llm_client.get_message_factory()._create_image_content_(
                document.identifier, document.data, detail
            )

    def create_pdf_content(self, document_type: FileUtilDocument, detail: str = "auto") -> list["ChatContent"]:
        config = self.llm_client.get_config()
        if not config:
            raise ValueError("LLMClientの設定が取得できませんでした。")

        use_custom = config.features.use_custom_pdf_analyzer
        if use_custom:
            return self.llm_client.get_message_factory()._create_custom_pdf_content_(
                document_type.identifier, document_type.data, detail=detail
                )
        else:
            return self.llm_client.get_message_factory()._create_pdf_content_(
                document_type.identifier, document_type.data, detail=detail)

    def create_pdf_content_from_file(self, file_path: str, detail: str = "auto") -> list["ChatContent"]:
            return self.create_pdf_content(FileUtilDocument.from_file(file_path), detail=detail)

    def create_pdf_content_from_url(self, file_url: str, detail: str = "auto") -> list["ChatContent"]:
        tmpdir = tempfile.TemporaryDirectory()
        atexit.register(tmpdir.cleanup)

        requests_verify, ca_bundle = self._get_network_download_options()
        file_paths = DownLoader.download_files(
            [WebRequestModel(url=file_url)],
            tmpdir.name,
            requests_verify=requests_verify,
            ca_bundle=ca_bundle,
        )
        return self.create_pdf_content_from_file(file_paths[0], detail=detail)

    def _create_custom_pdf_content_from_file(self, file_path: str, detail: str = "auto") -> list["ChatContent"]:
        '''
        PDFファイルのバイトデータから、テキスト抽出と画像抽出を行い、ChatContentのリストを生成して返す
        '''
        with open(file_path, "rb") as pdf_file:
            document_type = FileUtilDocument(data=pdf_file.read(), identifier=file_path)
        return self.llm_client.get_message_factory()._create_custom_pdf_content_(
            document_type.identifier, document_type.data, detail=detail
        )

    def create_office_content(self, document_type: FileUtilDocument, detail: str) -> list["ChatContent"]:
        '''
        複数のOfficeドキュメントとプロンプトからドキュメント解析を行う。各ドキュメントのテキスト抽出、各ドキュメントの説明、プロンプト応答を生成して返す
        '''
        effective_config = self._get_effective_config()
        office2pdf_method = effective_config.office2pdf.method
        if not Office2PDFUtil.is_conversion_available(config=effective_config):
            logger.info(
                "Office2PDF method %s is unavailable. Falling back to direct office text extraction for %s",
                office2pdf_method,
                document_type.identifier,
            )
            return self._create_office_text_fallback_content(document_type)

        try:
            # Officeドキュメントを一時的にPDFに変換する
            temp_dir = tempfile.TemporaryDirectory()
            atexit.register(temp_dir.cleanup)
            temp_file_path = os.path.join(
                temp_dir.name, 
                f"{os.path.basename(document_type.identifier)}_{uuid.uuid4()}.pdf"
                )
            Office2PDFUtil.create_pdf_from_document_bytes(
                input_bytes=document_type.data,
                output_path=temp_file_path,
                config=effective_config,
            )
            # 元ファイルからPDFに変換した旨の説明を追加
            explanation_text = f"""
                {temp_file_path}は、元のOfficeドキュメント: {document_type.identifier} をPDFに変換したものです。
                ユーザーにどのファイルを元にしたのかを伝えるため、回答を行う際には、元のファイル名を使用してください。
                """
            explanation_content = self.create_text_content(text=explanation_text)
            pdf_contents = [explanation_content]

            pdf_contents.extend(self.create_pdf_content_from_file(temp_file_path, detail=detail))

            return pdf_contents
        except Exception:
            logger.warning(
                "Failed to convert office document to PDF with method %s. Falling back to direct text extraction for %s",
                office2pdf_method,
                document_type.identifier,
                exc_info=True,
            )
            return self._create_office_text_fallback_content(document_type)

    def create_office_content_from_file(
            self, file_path: str, detail: str = "auto"
            ) -> list["ChatContent"]:
        '''
        複数のOfficeドキュメントとプロンプトからドキュメント解析を行う。各ドキュメントのテキスト抽出、各ドキュメントの説明、プロンプト応答を生成して返す
        '''
        with open(file_path, "rb") as office_file:
            document_type = FileUtilDocument(data=office_file.read(), identifier=file_path)
        return self.create_office_content(document_type, detail=detail)

    def create_office_content_from_url(self, file_url: str, detail: str = "auto") -> list["ChatContent"]:
        '''
        複数のOfficeドキュメントとプロンプトからドキュメント解析を行う。各ドキュメントのテキスト抽出、各ドキュメントの説明、プロンプト応答を生成して返す
        '''
        tmpdir = tempfile.TemporaryDirectory()
        atexit.register(tmpdir.cleanup)

        requests_verify, ca_bundle = self._get_network_download_options()
        file_paths = DownLoader.download_files(
            [WebRequestModel(url=file_url)],
            tmpdir.name,
            requests_verify=requests_verify,
            ca_bundle=ca_bundle,
        )

        office_contents = []
        for file_path in file_paths:
            content = self.create_office_content_from_file(
                file_path, detail=detail)
            office_contents.extend(content)

        return office_contents

    def create_multi_format_content(
            self, document_type: FileUtilDocument, detail: str = "auto"
            ) -> list["ChatContent"]:
        '''
        複数形式ファイルから、テキスト抽出と画像抽出を行い、ChatContentのリストを生成して返す
        '''

        if document_type.is_text():
            text = document_type.data.decode('utf-8')
            return [self.create_text_content(text)]

        if document_type.is_image():
            return self.create_image_content(document_type, detail)

        if document_type.is_pdf():
            return self.create_pdf_content(document_type, detail=detail)

        if document_type.is_office_document():
            return self.create_office_content(document_type, detail=detail)

        raise ValueError(f"Unsupported document type for file: {document_type.identifier}")

    def create_multi_format_contents_from_file(
            self, file_path: str, detail: str = "auto"
            ) -> list["ChatContent"]:
        '''
        複数形式ファイルから、テキスト抽出と画像抽出を行い、ChatContentのリストを生成して返す
        '''

        document_type = FileUtilDocument.from_file(document_path=file_path)

        if document_type.is_text():
            with open(file_path, "r", encoding="utf-8") as text_file:
                text_data = text_file.read()
            return [self.create_text_content(text_data)]

        if document_type.is_image():
            return self.create_image_content_from_file(file_path, detail)

        if document_type.is_pdf():
            return self.create_pdf_content_from_file(file_path, detail=detail)

        if document_type.is_office_document():
            return self.create_office_content_from_file(file_path, detail=detail)

        raise ValueError(f"Unsupported document type for file: {file_path}")

    def create_multi_format_contents_from_url(
            self, file_url: str, detail: str = "auto"
            ) -> list["ChatContent"]:
        '''
        複数形式ファイルから、テキスト抽出と画像抽出を行い、ChatContentのリストを生成して返す
        '''
        tmpdir = tempfile.TemporaryDirectory()
        atexit.register(tmpdir.cleanup)
        requests_verify, ca_bundle = self._get_network_download_options()
        file_paths = DownLoader.download_files(
            [WebRequestModel(url=file_url)],
            tmpdir.name,
            requests_verify=requests_verify,
            ca_bundle=ca_bundle,
        )

        return self.create_multi_format_contents_from_file(
            file_paths[0], detail=detail
        )

