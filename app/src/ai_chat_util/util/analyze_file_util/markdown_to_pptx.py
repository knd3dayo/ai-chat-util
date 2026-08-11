from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from .markdown_parser import MarkdownParser


class MarkdownToPptxUtil:
    @classmethod
    def convert_markdown_to_pptx(cls, markdown_text: str, output_path: str | Path) -> str:
        prs = Presentation()
        blocks = MarkdownParser.parse(markdown_text)
        slides: list[list[str]] = []
        current_title: str | None = None
        current_body: list[str] = []

        for block in blocks:
            if block.type == "heading" and block.level <= 2:
                if current_title is not None or current_body:
                    slides.append([current_title or f"Slide {len(slides) + 1}", *current_body])
                current_title = block.text
                current_body = []
            elif block.type == "paragraph" and block.text:
                current_body.append(block.text)
            elif block.type == "list_item" and block.text:
                current_body.append(f"• {block.text}")
            elif block.type == "code_block" and block.code:
                current_body.append(block.code)
            elif block.type == "table" and block.rows:
                current_body.append(" | ".join(block.rows[0]))
            elif block.type == "image" and block.image_path:
                current_body.append(f"[image] {block.image_path}")

        if current_title is not None or current_body:
            slides.append([current_title or f"Slide {len(slides) + 1}", *current_body])

        if not slides:
            slides = [["Untitled", ""]]

        for index, block in enumerate(slides):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            body = slide.placeholders[1]
            title.text = block[0] if block else f"Slide {index + 1}"
            content = "\n".join(block[1:]) if len(block) > 1 else ""
            text_frame = body.text_frame
            text_frame.clear()
            paragraphs = text_frame.paragraphs
            paragraph = paragraphs[0]
            paragraph.alignment = PP_ALIGN.LEFT
            run = paragraph.add_run()
            run.text = content
            run.font.size = Pt(18)
            if "[image]" in content:
                image_path = content.split("[image] ", 1)[1]
                try:
                    slide.shapes.add_picture(str(Path(image_path)), Inches(1.2), Inches(1.5), width=Inches(4.5))
                except FileNotFoundError:
                    pass

        output = Path(output_path)
        output.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output))
        return str(output)
